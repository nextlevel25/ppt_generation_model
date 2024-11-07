import datetime
import json
import os
import re

import openai
import requests
import telebot
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
import ssl
ssl._create_default_https_context = ssl._create_unverified_context


# Set up OpenAI API key and Telegram Bot token
openai.api_key = "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"
BOT_TOKEN = "6772673902:AAGg4zzRXMslhqn9CA0O2uFzePJiekuQXww"
GROUP_CHAT_ID = '-1002091120118'

# Initialize the Telegram bot
bot = telebot.TeleBot(BOT_TOKEN)

# Dictionary to store accumulated question data
new_data = {
    "Question": [],
    "Options": [],
    "Answer": [],
    "Explanation": []
}

# Define color and font settings for PPT
BACKGROUND_COLOR = RGBColor(250, 250, 250)
BORDER_COLOR = RGBColor(29, 115, 150)
QUESTION_TEXT_COLOR = RGBColor(0, 0, 0)
OPTION_TEXT_COLOR = RGBColor(0, 112, 192)
ANSWER_TEXT_COLOR = RGBColor(139, 0, 0)
EXPLANATION_TEXT_COLOR = RGBColor(139, 0, 10)
TEXT_SIZE = Pt(34)
FONT_NAME = "Verdana"

# Function to clean and extract JSON from the API response
def clean_json_response(formatted_data):
    """
    Cleans and extracts valid JSON array content from the response.
    """
    formatted_data = formatted_data.strip("`")
    json_match = re.search(r'\[\s*{\s*".*?}\s*]\s*$', formatted_data, re.DOTALL)

    if json_match:
        json_content = json_match.group(0)
        return json_content
    else:
        return None

# Modified OpenAI API function
def get_single_question_data(extracted_text, max_retries=3, max_tokens=500):
    """
    Queries ChatGPT to retrieve two questions at a time from the extracted text.
    """
    all_data = {
        "Question": [],
        "Options": [],
        "Answer": [],
        "Explanation": []
    }
    offset = 0

    while offset < len(extracted_text):
        prompt = (
            "Extract exactly two questions from the following OCR text, starting from the specified offset position. "
            "Return only the extracted questions and their details in a JSON array. Strictly adhere to the exact format specified below, "
            "and do not include any additional text, notes, interpretations, or unrelated content. Only output structured JSON in this format:\n\n"
            "[\n"
            "  {\n"
            "    \"Question\": \"Complete question text here as it appears in the OCR content.\",\n"
            "    \"Options\": [\"Option A text\", \"Option B text\", \"Option C text\", \"Option D text\"],\n"
            "    \"Answer\": \"Correct answer text or choice label here (e.g., 'Option B').\",\n"
            "    \"Explanation\": \"Explanation text here for the correct answer.\"\n"
            "  },\n"
            "  {\n"
            "    \"Question\": \"Next complete question text here as it appears in the OCR content.\",\n"
            "    \"Options\": [\"Option A text\", \"Option B text\", \"Option C text\", \"Option D text\"],\n"
            "    \"Answer\": \"Correct answer text or choice label here (e.g., 'Option C').\",\n"
            "    \"Explanation\": \"Explanation text here for the correct answer.\"\n"
            "  }\n"
            "]\n\n"
            "Make sure each question includes:\n"
            "- The full question text under the \"Question\" key.\n"
            "- Exactly four options under the \"Options\" key, listed in order as they appear, each option enclosed in quotation marks and separated by commas.\n"
            "- The correct answer under the \"Answer\" key, as provided in the OCR text, either as a label (e.g., 'Option B') or the full option text.\n"
            "- The explanation under the \"Explanation\" key, if available, otherwise leave it blank (e.g., \"Explanation\": \"\").\n\n"
            "If the OCR content does not include all required details for each question, do your best to provide placeholders (e.g., \"Explanation\": \"\") "
            "but do not omit any keys.\n\n"
            f"Here is the text starting from position {offset}:\n\n{extracted_text[offset:]}"
        )

        attempts = 0
        while attempts < max_retries:
            try:
                response = openai.ChatCompletion.create(
                    model="chatgpt-4o-latest",
                    messages=[
                        {"role": "system", "content": "You are a helpful assistant that strictly returns JSON."},
                        {"role": "user", "content": prompt}
                    ],
                    max_tokens=max_tokens
                )

                formatted_data = response.choices[0].message['content']
                print(f"Raw API response for offset {offset}:\n{formatted_data}")

                # Apply the improved cleaner function
                cleaned_data = clean_json_response(formatted_data)
                if cleaned_data is None:
                    raise ValueError("No JSON object detected.")

                single_question_data = json.loads(cleaned_data)

                # Append question data to all_data
                for question_data in single_question_data:
                    all_data["Question"].append(question_data["Question"])
                    all_data["Options"].append(question_data["Options"])
                    all_data["Answer"].append(question_data["Answer"])
                    all_data["Explanation"].append(question_data["Explanation"])

                # Update offset based on processed text length
                offset += len(cleaned_data)
                break

            except (json.JSONDecodeError, KeyError, Exception) as e:
                print(f"Error at offset {offset}, attempt {attempts + 1}: {e}")
                attempts += 1

        if attempts == max_retries:
            print(f"Failed after {max_retries} attempts at offset {offset}. Exiting.")
            break

    return all_data

# Accumulate questions from text extracted from an image
def accumulate_questions_from_text(extracted_text):
    questions_data = get_single_question_data(extracted_text)
    new_data["Question"].extend(questions_data["Question"])
    new_data["Options"].extend(questions_data["Options"])
    new_data["Answer"].extend(questions_data["Answer"])
    new_data["Explanation"].extend(questions_data["Explanation"])

# Function to extract text from an image using OCR
import cv2
from PIL import Image
import pytesseract

def preprocess_image(image_path):
    image = cv2.imread(image_path, cv2.IMREAD_GRAYSCALE)
    _, thresh_image = cv2.threshold(image, 128, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    denoised_image = cv2.fastNlMeansDenoising(thresh_image, None, 30, 7, 21)
    return Image.fromarray(denoised_image)

import easyocr


def extract_text_from_image(file_id):
    file_info = bot.get_file(file_id)
    file_url = f"https://api.telegram.org/file/bot{BOT_TOKEN}/{file_info.file_path}"
    file_path = f"temp_image_{datetime.datetime.now().strftime('%Y%m%d%H%M%S')}.jpg"

    # Download the file
    with open(file_path, "wb") as f:
        f.write(requests.get(file_url).content)

    # Check if file was downloaded successfully
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File could not be downloaded: {file_path}")

    try:
        # Process the image with EasyOCR
        reader = easyocr.Reader(['en', 'hi'])  # Specify languages needed
        result = reader.readtext(file_path, detail=0)
        extracted_text = "\n".join(result)
        print("text %s", extracted_text)
    finally:
        # Clean up: Remove the temporary image file
        os.remove(file_path)

    return extracted_text

# Function to create PowerPoint from accumulated data
def create_presentation_from_data(data):
    """
    Creates a PowerPoint presentation from the structured question data.
    Adds slides for each question and its answer with explanations.
    """
    prs = Presentation()
    prs.slide_width = Inches(14)
    prs.slide_height = Inches(7.5)
    title = "Quiz Questions"

    for i, question in enumerate(data["Question"]):
        options = data["Options"][i]
        answer = data["Answer"][i]
        explanation = data["Explanation"][i]

        add_slide(prs, title, i + 1, question, options, answer)
        add_slide(prs, title, i + 1, question, options, answer, is_answer_slide=True, explanation=explanation)

    prs.save("formatted_quiz.pptx")

# Add slide function
def add_slide(prs, title, question_number, question, options, answer=None, is_answer_slide=False, explanation=None):
    """
    Adds a slide with specified content (question and options or answer and explanation).
    If is_answer_slide is True, displays only the answer and explanation.
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set the background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR

    # Add the left-side color border
    border_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(1.5), Inches(7.5)
    )
    border_box.fill.solid()
    border_box.fill.fore_color.rgb = BORDER_COLOR
    border_box.line.fill.background()

    # Add the heading
    heading_box = slide.shapes.add_textbox(Inches(2), Inches(0.2), Inches(12), Inches(1))
    heading_frame = heading_box.text_frame
    heading_frame.text = title
    heading_paragraph = heading_frame.paragraphs[0]
    heading_paragraph.font.size = Pt(32)
    heading_paragraph.font.bold = True
    heading_paragraph.font.color.rgb = BORDER_COLOR
    heading_paragraph.alignment = PP_ALIGN.CENTER
    heading_paragraph.font.name = FONT_NAME

    if not is_answer_slide:
        # Add question text
        question_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(11.5), Inches(3))
        question_frame = question_box.text_frame
        question_frame.text = f"{question_number}. {question}"
        question_paragraph = question_frame.paragraphs[0]
        question_paragraph.font.size = TEXT_SIZE
        question_paragraph.font.bold = True
        question_paragraph.font.color.rgb = QUESTION_TEXT_COLOR
        question_paragraph.alignment = PP_ALIGN.LEFT
        question_paragraph.font.name = FONT_NAME

        # Add options
        options_top_margin = 4
        option_box = slide.shapes.add_textbox(Inches(2.2), Inches(options_top_margin), Inches(11.5), Inches(2.5))
        option_frame = option_box.text_frame
        option_frame.word_wrap = True
        option_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for option in options:
            p = option_frame.add_paragraph()
            p.text = option
            p.font.size = TEXT_SIZE
            p.font.color.rgb = OPTION_TEXT_COLOR
            p.alignment = PP_ALIGN.LEFT
            p.font.name = FONT_NAME

    else:
        # Add correct answer
        answer_box = slide.shapes.add_textbox(Inches(2), Inches(1.7), Inches(11.5), Inches(1))
        answer_frame = answer_box.text_frame
        answer_frame.text = f"Correct Answer: {answer}"
        answer_paragraph = answer_frame.paragraphs[0]
        answer_paragraph.font.size = TEXT_SIZE
        answer_paragraph.font.bold = True
        answer_paragraph.font.color.rgb = ANSWER_TEXT_COLOR
        answer_paragraph.alignment = PP_ALIGN.LEFT
        answer_paragraph.font.name = FONT_NAME

        # Add explanation
        explanation_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(10.5), Inches(2.5))
        explanation_frame = explanation_box.text_frame
        explanation_frame.text = "Explanation: " + explanation
        explanation_paragraph = explanation_frame.paragraphs[0]
        explanation_paragraph.font.size = TEXT_SIZE - 4
        explanation_paragraph.font.color.rgb = EXPLANATION_TEXT_COLOR
        explanation_paragraph.alignment = PP_ALIGN.LEFT
        explanation_paragraph.font.name = FONT_NAME

# Telegram handler for receiving images
@bot.message_handler(content_types=['photo'])
def handle_image_message(message):
    photo = message.photo[-1]
    file_id = photo.file_id
    extracted_text = extract_text_from_image(file_id)

    if extracted_text:
        accumulate_questions_from_text(extracted_text)
        bot.send_message(message.chat.id, "Questions added. Send more images or type 'nextlevel' for PPT.")
    else:
        bot.send_message(message.chat.id, "No text detected in the image.")

@bot.message_handler(func=lambda message: message.text.lower() == "nextlevel")
def handle_nextlevel_command(message):
    if new_data["Question"]:
        # Generate PPTX
        create_presentation_from_data(new_data)
        pptx_file_path = "formatted_quiz.pptx"
        with open(pptx_file_path, "rb") as ppt_file:
            bot.send_document(message.chat.id, ppt_file)
            print("PPT file sent")

        # Clear accumulated data
        for key in new_data:
            new_data[key].clear()
    else:
        bot.send_message(message.chat.id, "No questions to create a presentation. Please send images first.")

# Start bot polling
bot.infinity_polling(timeout=60, long_polling_timeout=60)

